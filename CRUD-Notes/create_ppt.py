from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
TITLE_BG = RGBColor(0x2C, 0x3E, 0x50)
ACCENT = RGBColor(0x34, 0x98, 0xDB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)
BULLET_TEXT = RGBColor(0x44, 0x44, 0x44)

slides_data = [
    {
        "number": "Assistant 1",
        "title": "Requirement Writer",
        "bullets": [
            "Reads and updates requirements in Confluence Notes App page.",
            "Ensures professional and neat formatting.",
            "Removes special characters from requirements.",
            "Provides update link after successful changes.",
        ],
    },
    {
        "number": "Assistant 2",
        "title": "Gap Analyser",
        "bullets": [
            "Acts as an experienced architect.",
            "Compares business requirements with application HTML.",
            "Identifies gaps and documents them in Confluence Gaps page.",
            "Content is prepared for user story creation by the next agent.",
        ],
    },
    {
        "number": "Assistant 3",
        "title": "User Story Generator",
        "bullets": [
            "Reads gaps from Confluence Gaps page.",
            "Generates Jira user stories with prefix 'codemie'.",
            "Updates the same Jira story each time (no new stories).",
            "Ensures stories are ready for development.",
        ],
    },
    {
        "number": "Assistant 4",
        "title": "Plan and Design",
        "bullets": [
            "Selects top 3 gaps for implementation.",
            "Creates a plan for these gaps in Confluence Plan page.",
            "Designs architecture flow charts for the gaps in Confluence Design page.",
            "Updates content in the same pages each time.",
        ],
    },
    {
        "number": "Assistant 5",
        "title": "Develop Assistant (Claude)",
        "bullets": [
            "Develops features based on gaps and user stories.",
            "Refers to attached files for implementation steps.",
            "Handles code deployment, PR raising, merging, and feature checks.",
        ],
    },
    {
        "number": "Assistant 6",
        "title": "Test",
        "bullets": [
            "Waits for user approval before starting tests.",
            "Writes tests for top three gaps using Playwright MCP.",
            "Documents manual tests in Confluence in simple English/Gherkin.",
            "Executes tests on http://localhost:3001/ or via HTML if not accessible.",
            "Saves Playwright report with screenshots and results in Confluence.",
            "Uses Playwright MCP for all testing tasks.",
        ],
    },
]


def add_bg_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def make_slide(data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Left accent bar
    add_bg_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT)

    # Top title band
    add_bg_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.6), TITLE_BG)

    # Number badge
    badge = add_bg_shape(slide, Inches(0.6), Inches(0.35), Inches(1.8), Inches(0.9), ACCENT)
    tf = badge.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data["number"]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Title text
    title_box = slide.shapes.add_textbox(Inches(2.8), Inches(0.3), Inches(9), Inches(1.0))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = data["title"]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Bullet content area
    content_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11), Inches(4.8))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(data["bullets"]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(14)
        p.space_before = Pt(6)

        # Bullet marker
        run_marker = p.add_run()
        run_marker.text = "\u25B8  "
        run_marker.font.size = Pt(20)
        run_marker.font.color.rgb = ACCENT
        run_marker.font.bold = True

        # Bullet text
        run_text = p.add_run()
        run_text.text = bullet
        run_text.font.size = Pt(20)
        run_text.font.color.rgb = BULLET_TEXT


for s in slides_data:
    make_slide(s)

output_path = os.path.join(os.path.dirname(__file__), "AI_Assistants_Pipeline.pptx")
prs.save(output_path)
print(f"PPT saved to: {output_path}")
